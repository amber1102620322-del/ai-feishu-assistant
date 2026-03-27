import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_structured_pptx(title: str, slides_data: list, output_path: str):
    """
    使用 python-pptx 创建一个美观的 PPT。
    slides_data: [{"title": "标题", "content": "内容"}, ...]
    """
    prs = Presentation()
    
    # 设置幻灯片大小为 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # 1. 标题页
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102) # 深蓝色

    subtitle = slide.placeholders[1]
    subtitle.text = "由飞书 AI 助手专业生成\n" + "-" * 20
    
    # 2. 内容页
    for data in slides_data:
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        
        # 标题风格
        title_shape = slide.shapes.title
        title_shape.text = data.get("title", "")
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.size = Pt(32)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0, 102, 204)

        # 正文风格
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.word_wrap = True
        
        content = data.get("content", "")
        # 如果内容是列表，按行拆分
        lines = content.split('\n')
        for i, line in enumerate(lines):
            line = line.strip()
            if not line: continue
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            # 去掉可能的 Markdown 符号
            clean_line = line.lstrip('- ').lstrip('* ').lstrip('123456789. ')
            p.text = clean_line
            p.level = 0
            p.font.size = Pt(20)

    # 3. 结束页
    end_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(end_slide_layout)
    slide.shapes.title.text = "感谢您的观看"
    slide.placeholders[1].text = "如有修改需求请随时联系 AI 助手"

    prs.save(output_path)
    return output_path
